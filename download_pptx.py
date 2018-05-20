"""

Created on May 17, 2018

Note: this script outputs a file, named with the date and a .pptx extension

@author: Rahul Shah


"""

import json
import calendar
import datetime as dt

# imports - don't remove any of these, even if they seem as if they aren't being used

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import *

# layout constants
TITLE = 0
TITLE_AND_CONTENT = 1
SECTION_HEADER = 2
TWO_CONTENT = 3
COMPARISON = 4
TITLE_ONLY = 5
BLANK = 6
CONTENT_WITH_CAPTION = 7
PICTURE_WITH_CAPTION = 8

# Parse data from json
with open('sample.json') as data_file:
    data = json.load(data_file)

    # Make a reference variable for the presentation
    prs = Presentation()

    # first slide
    first_slide = prs.slide_layouts[BLANK]
    slide1 = prs.slides.add_slide(first_slide)
    left = top = width = height = Inches(1.0)
    shapes = slide1.shapes

    # shape1 = shapes.add_shape(
    #     MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    # )

    l1 = Inches(1)
    t1 = Inches(1)
    w1 = Inches(3)
    h1 = Inches(1)

    # made for schedule table
    shape1 = slide1.shapes

    txBox = slide1.shapes.add_textbox(l1, t1, w1, h1)
    tf = txBox.text_frame
    tf.word_wrap = True

    # date stuff
    p = tf.add_paragraph()
    days = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]

    unformatted_date = ((data["board"]["timestamp"])[0:10])
    d = dt.datetime.strptime(unformatted_date, "%Y-%m-%d")
    day = days[d.weekday()] 
    month = calendar.month_name[int(unformatted_date[5:7])]
    day_of_month = unformatted_date[8:10]
    year = unformatted_date[0:4]

    full_date = day + ', ' + month + ' ' + day_of_month + ', ' + year
    school_date = 'Day: ' + data["board"]["day"]
    formatted_date = full_date + '\n\n' + school_date

    p.text = formatted_date

    p.font.bold = False

    p = tf.add_paragraph()
    p.text = '\n\n' + "Announcements: " + '\n\n' + data["board"]["announcements"]
    p.font.size = Pt(24)

    p = tf.add_paragraph()
    p.text = '\n\n' + "Quote: " + data["board"]["quote"]
    p.font.size = Pt(18)
    p.font.bold = True

    # optional font color text

    # text_frame = shape1.text_frame
    # text_frame.text = data["board"]["announcements"]
    # text_frame.margin_bottom = Inches(-1.08)
    # text_frame.margin_left = 0
    # text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    # text_frame.word_wrap = True
    # text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    rows = 8
    cols = 3
    t_left = Inches(4.55)
    t_top = Inches(0.0)
    t_width = Inches(5.90)
    t_height = Inches(7.50)

    table = shape1.add_table(rows, cols, t_left, t_top, t_width, t_height).table

    # set column widths
    table.columns[0].width = Inches(1.75)
    table.columns[1].width = Inches(1.75)
    table.columns[1].width = Inches(1.75)

    # set column heights
    table.columns[0].height = Inches(2.25)
    table.columns[1].height = Inches(2.25)
    table.columns[2].height = Inches(2.25)

    # Don't change this (table headers)
    period = table.cell(0, 0).text = 'Period'
    start = table.cell(0, 1).text = 'Start'
    end = table.cell(0, 2).text = 'End'

    period.center(1)
    start.center(1)
    end.center(1)

    # period color
    cell1 = table.rows[0].cells[0]
    paragraph1 = cell1.text_frame.paragraphs[0]
    paragraph1.font.color.rgb = RGBColor(00, 255, 00)

    # start color
    cell1 = table.rows[0].cells[1]
    paragraph1 = cell1.text_frame.paragraphs[0]
    paragraph1.font.color.rgb = RGBColor(00, 255, 00)

    # end color
    cell1 = table.rows[0].cells[2]
    paragraph1 = cell1.text_frame.paragraphs[0]
    paragraph1.font.color.rgb = RGBColor(00, 255, 00)

    # Put period numbers
    for index, item in enumerate(data["schedule"]):
        table.cell((index + 1), 0).text = item["number"]

    # Put start time of class
    for index, item in enumerate(data["schedule"]):
        table.cell((index + 1), 1).text = item["start"]

    # Put end time of class
    for index, item in enumerate(data["schedule"]):
        table.cell((index + 1), 2).text = item["end"]

    # for non-school days
    if data["board"]["message_only"]:
        # slide 2
        title_only_slide_layout = prs.slide_layouts[TITLE_ONLY]
        a_slide = prs.slides.add_slide(title_only_slide_layout)
        l2 = Inches(1)
        t2 = Inches(0)
        w2 = Inches(1)
        h2 = Inches(1)
        txBox = a_slide.shapes.add_textbox(l2, t2, w2, h2)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = 'No School!'
        p.font.size = Pt(36)
        p.font.bold = True

        shape2 = a_slide.shapes

    elif not data["board"]["message_only"]:
        # slide 2
        title_only_slide_layout = prs.slide_layouts[TITLE_ONLY]
        slide2 = prs.slides.add_slide(title_only_slide_layout)

        # Header text
        title = slide2.shapes.title

        title.text = "Absent Teachers"

        # Create a SlideShape object to later create a table from
        shape2 = slide2.shapes

        # dimensions of table
        rows = len(data["absents"][0]["teacher"]) - 2
        cols = 2
        left = Inches(0.0)
        top = Inches(1.5)
        width = Inches(15.0)
        height = Inches(0.8)

        table = shape2.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(5.5)
        table.columns[1].width = Inches(4.5)

        # Don't change this (table header)
        table.cell(0, 0).text = 'Teacher'
        table.cell(0, 1).text = 'Hour'

        # Put teacher names
        for index, item in enumerate(data["absents"]):
            table.cell((index + 1), 0).text = item["teacher"]

        # Put teacher hours
        for index, item in enumerate(data["absents"]):
            table.cell((index + 1), 1).text = item["hours"]

# Generate slides
name = unformatted_date + '.pptx'
prs.save(name)

print('Successfully made powerpoint as: ' + name)  # remove this
