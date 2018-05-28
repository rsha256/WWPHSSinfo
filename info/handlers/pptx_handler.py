import datetime
import calendar
import textwrap

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_pptx(data):

    if data['board']['message_only']:
        return _create_message_only(data)
    else:
        return _create_full(data)


def _create_full(data):
    date_obj = datetime.datetime.strptime(data['board']['timestamp'].split('+')[0], '%Y-%m-%d %H:%M:%S.%f').date()
    formatted_date = "{}, {:%b %d, %Y}".format(calendar.day_name[date_obj.weekday()], date_obj)

    prs = Presentation()

    # ### Title Slide ### #
    title_slide_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(title_slide_layout)

    # # Date # #
    date_text_box = title_slide.shapes.add_textbox(left=Inches(0.2), top=Inches(0), width=Inches(10), height=Inches(1))
    date_text_frame = date_text_box.text_frame

    date_p = date_text_frame.add_paragraph()

    date_p.text = formatted_date
    date_p.font.color.rgb = RGBColor(0, 100, 0)
    date_p.font.name = 'Roberto Light'
    date_p.font.size = Pt(36)

    # Content
    date_content_text_box = title_slide.shapes.add_textbox(left=Inches(1), top=Inches(1), width=Inches(0.5), height=Inches(5))
    date_content_frame = date_text_box.text_frame

    date_p = date_content_frame.add_paragraph()
    date_p.word_wrap = True
    date_p.text = "Day: " + data["board"]["day"]
    date_p.font.name = 'Roberto Light'
    date_p.font.size = Pt(20)

    # # Announcements # #

    # Heading
    announcements_heading_text_box = title_slide.shapes.add_textbox(left=Inches(0.5), top=Inches(0.75), width=Inches(10), height=Inches(1))
    announcements_heading_frame = announcements_heading_text_box.text_frame

    announcements_heading_p = announcements_heading_frame.add_paragraph()
    announcements_heading_p.text = "\n\n" + "Announcements:"
    announcements_heading_p.font.color.rgb = RGBColor(0, 100, 0)
    announcements_heading_p.font.name = 'Roberto Light'
    announcements_heading_p.font.size = Pt(24)

    # Content
    announcements_content_text_box = title_slide.shapes.add_textbox(left=Inches(1.1), top=Inches(1), width=Inches(0.5), height=Inches(5))
    announcements_content_frame = announcements_content_text_box.text_frame

    announcements_content_p = announcements_content_frame.add_paragraph()
    announcements_content_p.word_wrap = True

    s = data['board']['announcements']
    a_text = '\n'.join(textwrap.wrap(s, 38))
    announcements_content_p.text = "\n\n\n" + a_text

    announcements_content_p.font.name = 'Roberto Light'
    announcements_content_p.font.size = Pt(20)

    # # Quote # #

    # Heading
    quote_heading_text_box = title_slide.shapes.add_textbox(left=Inches(0.5), top=Inches(5.5), width=Inches(5), height=Inches(1))
    quote_heading_frame = quote_heading_text_box.text_frame

    quote_heading_frame_p = quote_heading_frame.add_paragraph()
    quote_heading_frame_p.text = "Quote:"
    quote_heading_frame_p.font.color.rgb = RGBColor(0, 100, 0)
    quote_heading_frame_p.font.name = 'Roberto Light'
    quote_heading_frame_p.font.size = Pt(24)

    # Content
    quote_content_text_box = title_slide.shapes.add_textbox(left=Inches(1.45), top=Inches(6), width=Inches(7), height=Inches(1))
    quote_content_frame = quote_content_text_box.text_frame

    quote_content_frame_p = quote_content_frame.add_paragraph()
    quote_content_frame_p.text = data['board']['quote']
    quote_content_frame_p.font.name = 'Roberto Light'
    quote_content_frame_p.font.size = Pt(20)
    quote_content_frame_p.word_wrap = True

    # # Schedule # #

    schedule_table = title_slide.shapes.add_table(rows=(len(data['schedule']) + 1), cols=3, left=Inches(5.5),
                                                  top=Inches(0),
                                                  width=Inches(4.5), height=Inches(5)).table

    schedule_table.cell(0, 0).text = "Period"
    schedule_table.cell(0, 1).text = "Start"
    schedule_table.cell(0, 2).text = "End"

    i = 1
    for s in data['schedule']:
        schedule_table.cell(i, 0).text = s['number']
        schedule_table.cell(i, 1).text = s['start']
        schedule_table.cell(i, 2).text = s['end']
        i += 1

    # ### Absent Teachers Slides ### #

    absent_lists = list(__chunks(data['absents'], 10))

    for chunk in absent_lists:
        absent_slide_layout = prs.slide_layouts[6]
        absent_slide = prs.slides.add_slide(absent_slide_layout)

        table = absent_slide.shapes.add_table(rows=11, cols=2, left=Inches(0), top=Inches(1), width=Inches(10),
                                              height=Inches(6)).table
        table.cell(0, 0).text = "Teacher"
        table.cell(0, 1).text = "Hours"

        j = 1
        for teacher in chunk:
            table.cell(j, 0).text = teacher['teacher']
            table.cell(j, 1).text = teacher['hours']
            j += 1

    return prs


def _create_message_only(data):
    date_obj = datetime.datetime.strptime(data['board']['timestamp'].split('+')[0], '%Y-%m-%d %H:%M:%S.%f').date()
    formatted_date = "{}, {:%b %d, %Y}".format(calendar.day_name[date_obj.weekday()], date_obj)

    prs = Presentation()

    # Slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    text_box = slide.shapes.add_textbox(left=Inches(0), top=Inches(0), width=Inches(10), height=Inches(1))
    text_frame = text_box.text_frame

    p = text_frame.add_paragraph()

    p.text = formatted_date
    p.font.color.rgb = RGBColor(0, 100, 0)
    p.font.name = 'Roberto Light'
    p.font.size = Pt(36)

    # Message
    text_box = slide.shapes.add_textbox(left=Inches(0), top=Inches(1), width=Inches(10), height=Inches(1))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()

    p.text = data['board']['message']

    p.font.name = 'Roberto Light'
    p.font.size = Pt(26)

    return prs


# Splits a list into evenly sized chunks, leaving leftovers
def __chunks(l, n):
    # For item i in a range that is a length of l,
    for i in range(0, len(l), n):
        # Create an index range for l of n items:
        yield l[i:i + n]